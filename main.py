from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches
import os
import cv2
from pytube import YouTube

my_secret = os.environ['HEROKU_API_KEY']
app = Flask(__name__)


def images_to_ppt(directory, ppt_file):
    # Create a new presentation
    prs = Presentation()

    # Get a list of all the image files in the directory
    images = [f for f in os.listdir(directory) if f.endswith(('.png', '.jpg', '.jpeg', '.gif'))]

    # Add a slide with an image for each image file
    for image_file in images:
        slide_layout = prs.slide_layouts[5]  # Use the blank slide layout (layout 5)
        slide = prs.slides.add_slide(slide_layout)

        # Define width and height of the image on the slide
        width =  Inches(10)
        height = Inches(7.5)
        left = 0
        top = 0

        # Add the image
        slide.shapes.add_picture(os.path.join(directory, image_file), left, top, width, height)

    # Save the presentation
    prs.save(ppt_file)

# Example usage:



def download_youtube_as_mp4(video_url, resolution):
    from pytube import YouTube

    # Create a YouTube object and get the video stream
    youtube = YouTube(video_url)
    video_stream = youtube.streams.filter(progressive=True, file_extension='mp4', res=resolution).first()

    if video_stream is None:
        print(f"No video found with resolution {resolution}")
        return None

    # Function to track the progress of the download
    def progress_function(stream, chunk, bytes_remaining):
        current = (stream.filesize - bytes_remaining)
        percent = (current / stream.filesize) * 100
        print(f"Downloaded {current} of {stream.filesize} bytes ({percent:.2f}%)")

    # Register the progress function
    youtube.register_on_progress_callback(progress_function)

    # Download the video stream as an MP4 file
    mp4_file_path = video_stream.download()

    return mp4_file_path

import cv2
import os

def extract_frames_cv2(video_path, seconds):
    vidcap = cv2.VideoCapture(video_path)
    fps = vidcap.get(cv2.CAP_PROP_FPS) # Gets the frames per second
    multiplier = fps * seconds

    # Create a new directory to store the frames
    directory_name = os.path.basename(video_path)[:21]
    if not os.path.exists(directory_name):
        os.makedirs(directory_name)
    frame_number = 0;
    while vidcap.isOpened():
        
        frame_id = int(round(vidcap.get(1)))  # Current frame number
        ret, frame = vidcap.read()

        if not ret:
            break

        if frame_id % multiplier == 0:
            # Save the frame in the new directory
            frame_number = frame_number + 1
            cv2.imwrite(os.path.join(directory_name, f"frame{frame_number}.jpg"), frame)

    vidcap.release()


@app.route('/download_and_convert', methods=['POST'])
def download_and_convert():
    video_url = request.json['video_url']
    resolution = request.json['resolution']
    seconds = request.json['seconds']
    ppt_file = 'output.pptx'

    extract_frames_cv2('/content/Building an Audio Transcription App with OpenAI Whisper and Streamlit.mp4', 30)
    images_to_ppt('/content/Building an Audio Tra', 'Building an Audio Transcription App3.pptx')

    

    return send_file(ppt_file, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
