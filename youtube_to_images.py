{
  "nbformat": 4,
  "nbformat_minor": 0,
  "metadata": {
    "colab": {
      "provenance": [],
      "authorship_tag": "ABX9TyMV32wsfwXoRtmaZgYGzggx",
      "include_colab_link": true
    },
    "kernelspec": {
      "name": "python3",
      "display_name": "Python 3"
    },
    "language_info": {
      "name": "python"
    }
  },
  "cells": [
    {
      "cell_type": "markdown",
      "metadata": {
        "id": "view-in-github",
        "colab_type": "text"
      },
      "source": [
        "<a href=\"https://colab.research.google.com/github/mileslilly2/youtube-to-slides/blob/master/youtube_to_images.py\" target=\"_parent\"><img src=\"https://colab.research.google.com/assets/colab-badge.svg\" alt=\"Open In Colab\"/></a>"
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "!pip install opencv-python\n",
        "!pip install moviepy\n",
        "!pip install imageio[ffmpeg]\n",
        "!pip install pydub\n",
        "!pip install pytube=12.0.0\n",
        "!pip install python-pptx\n"
      ],
      "metadata": {
        "colab": {
          "base_uri": "https://localhost:8080/"
        },
        "id": "tdH1cX2oro7P",
        "outputId": "949d0fc0-5f30-4489-ffd0-faea6c7b1b9a"
      },
      "execution_count": 1,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stdout",
          "text": [
            "Requirement already satisfied: opencv-python in /usr/local/lib/python3.10/dist-packages (4.8.0.76)\n",
            "Requirement already satisfied: numpy>=1.21.2 in /usr/local/lib/python3.10/dist-packages (from opencv-python) (1.23.5)\n",
            "Requirement already satisfied: moviepy in /usr/local/lib/python3.10/dist-packages (1.0.3)\n",
            "Requirement already satisfied: decorator<5.0,>=4.0.2 in /usr/local/lib/python3.10/dist-packages (from moviepy) (4.4.2)\n",
            "Requirement already satisfied: tqdm<5.0,>=4.11.2 in /usr/local/lib/python3.10/dist-packages (from moviepy) (4.66.1)\n",
            "Requirement already satisfied: requests<3.0,>=2.8.1 in /usr/local/lib/python3.10/dist-packages (from moviepy) (2.31.0)\n",
            "Requirement already satisfied: proglog<=1.0.0 in /usr/local/lib/python3.10/dist-packages (from moviepy) (0.1.10)\n",
            "Requirement already satisfied: numpy>=1.17.3 in /usr/local/lib/python3.10/dist-packages (from moviepy) (1.23.5)\n",
            "Requirement already satisfied: imageio<3.0,>=2.5 in /usr/local/lib/python3.10/dist-packages (from moviepy) (2.31.6)\n",
            "Requirement already satisfied: imageio-ffmpeg>=0.2.0 in /usr/local/lib/python3.10/dist-packages (from moviepy) (0.4.9)\n",
            "Requirement already satisfied: pillow<10.1.0,>=8.3.2 in /usr/local/lib/python3.10/dist-packages (from imageio<3.0,>=2.5->moviepy) (9.4.0)\n",
            "Requirement already satisfied: setuptools in /usr/local/lib/python3.10/dist-packages (from imageio-ffmpeg>=0.2.0->moviepy) (67.7.2)\n",
            "Requirement already satisfied: charset-normalizer<4,>=2 in /usr/local/lib/python3.10/dist-packages (from requests<3.0,>=2.8.1->moviepy) (3.3.2)\n",
            "Requirement already satisfied: idna<4,>=2.5 in /usr/local/lib/python3.10/dist-packages (from requests<3.0,>=2.8.1->moviepy) (3.6)\n",
            "Requirement already satisfied: urllib3<3,>=1.21.1 in /usr/local/lib/python3.10/dist-packages (from requests<3.0,>=2.8.1->moviepy) (2.0.7)\n",
            "Requirement already satisfied: certifi>=2017.4.17 in /usr/local/lib/python3.10/dist-packages (from requests<3.0,>=2.8.1->moviepy) (2023.11.17)\n",
            "Requirement already satisfied: imageio[ffmpeg] in /usr/local/lib/python3.10/dist-packages (2.31.6)\n",
            "Requirement already satisfied: numpy in /usr/local/lib/python3.10/dist-packages (from imageio[ffmpeg]) (1.23.5)\n",
            "Requirement already satisfied: pillow<10.1.0,>=8.3.2 in /usr/local/lib/python3.10/dist-packages (from imageio[ffmpeg]) (9.4.0)\n",
            "Requirement already satisfied: imageio-ffmpeg in /usr/local/lib/python3.10/dist-packages (from imageio[ffmpeg]) (0.4.9)\n",
            "Requirement already satisfied: psutil in /usr/local/lib/python3.10/dist-packages (from imageio[ffmpeg]) (5.9.5)\n",
            "Requirement already satisfied: setuptools in /usr/local/lib/python3.10/dist-packages (from imageio-ffmpeg->imageio[ffmpeg]) (67.7.2)\n",
            "Collecting pydub\n",
            "  Downloading pydub-0.25.1-py2.py3-none-any.whl (32 kB)\n",
            "Installing collected packages: pydub\n",
            "Successfully installed pydub-0.25.1\n",
            "\u001b[31mERROR: Invalid requirement: 'pytube=12.0.0'\n",
            "Hint: = is not a valid operator. Did you mean == ?\u001b[0m\u001b[31m\n",
            "\u001b[0mCollecting python-pptx\n",
            "  Downloading python_pptx-0.6.23-py3-none-any.whl (471 kB)\n",
            "\u001b[2K     \u001b[90m━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\u001b[0m \u001b[32m471.6/471.6 kB\u001b[0m \u001b[31m2.7 MB/s\u001b[0m eta \u001b[36m0:00:00\u001b[0m\n",
            "\u001b[?25hRequirement already satisfied: lxml>=3.1.0 in /usr/local/lib/python3.10/dist-packages (from python-pptx) (4.9.3)\n",
            "Requirement already satisfied: Pillow>=3.3.2 in /usr/local/lib/python3.10/dist-packages (from python-pptx) (9.4.0)\n",
            "Collecting XlsxWriter>=0.5.7 (from python-pptx)\n",
            "  Downloading XlsxWriter-3.1.9-py3-none-any.whl (154 kB)\n",
            "\u001b[2K     \u001b[90m━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\u001b[0m \u001b[32m154.8/154.8 kB\u001b[0m \u001b[31m16.7 MB/s\u001b[0m eta \u001b[36m0:00:00\u001b[0m\n",
            "\u001b[?25hInstalling collected packages: XlsxWriter, python-pptx\n",
            "Successfully installed XlsxWriter-3.1.9 python-pptx-0.6.23\n"
          ]
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "!pip install --upgrade pytube\n"
      ],
      "metadata": {
        "colab": {
          "base_uri": "https://localhost:8080/"
        },
        "id": "W4-lAVepivnr",
        "outputId": "f750ee17-042f-41ae-e134-2c66936fb369"
      },
      "execution_count": 2,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stdout",
          "text": [
            "Collecting pytube\n",
            "  Downloading pytube-15.0.0-py3-none-any.whl (57 kB)\n",
            "\u001b[?25l     \u001b[90m━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\u001b[0m \u001b[32m0.0/57.6 kB\u001b[0m \u001b[31m?\u001b[0m eta \u001b[36m-:--:--\u001b[0m\r\u001b[2K     \u001b[91m━━━━━━━━━━━━━━━━━━━━━\u001b[0m\u001b[90m╺\u001b[0m\u001b[90m━━━━━━━━━━━━━━━━━━\u001b[0m \u001b[32m30.7/57.6 kB\u001b[0m \u001b[31m635.5 kB/s\u001b[0m eta \u001b[36m0:00:01\u001b[0m\r\u001b[2K     \u001b[90m━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\u001b[0m \u001b[32m57.6/57.6 kB\u001b[0m \u001b[31m763.5 kB/s\u001b[0m eta \u001b[36m0:00:00\u001b[0m\n",
            "\u001b[?25hInstalling collected packages: pytube\n",
            "Successfully installed pytube-15.0.0\n"
          ]
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [],
      "metadata": {
        "id": "Z0RuE24jL5gn"
      },
      "execution_count": 2,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "from pytube import YouTube\n",
        "import cv2\n",
        "import os\n",
        "import math\n",
        "from pptx import Presentation\n",
        "from pptx.util import Inches\n",
        "import os\n",
        "from google.colab import drive\n",
        "drive.mount('/content/drive')\n",
        "\n"
      ],
      "metadata": {
        "id": "QjgZRxkXsBZ3",
        "colab": {
          "base_uri": "https://localhost:8080/",
          "height": 394
        },
        "outputId": "404093d6-4441-4232-eb34-03c5775e57ac"
      },
      "execution_count": 1,
      "outputs": [
        {
          "output_type": "error",
          "ename": "ModuleNotFoundError",
          "evalue": "ignored",
          "traceback": [
            "\u001b[0;31m---------------------------------------------------------------------------\u001b[0m",
            "\u001b[0;31mModuleNotFoundError\u001b[0m                       Traceback (most recent call last)",
            "\u001b[0;32m<ipython-input-1-ceb4ad0ca5ba>\u001b[0m in \u001b[0;36m<cell line: 1>\u001b[0;34m()\u001b[0m\n\u001b[0;32m----> 1\u001b[0;31m \u001b[0;32mfrom\u001b[0m \u001b[0mpytube\u001b[0m \u001b[0;32mimport\u001b[0m \u001b[0mYouTube\u001b[0m\u001b[0;34m\u001b[0m\u001b[0;34m\u001b[0m\u001b[0m\n\u001b[0m\u001b[1;32m      2\u001b[0m \u001b[0;32mimport\u001b[0m \u001b[0mcv2\u001b[0m\u001b[0;34m\u001b[0m\u001b[0;34m\u001b[0m\u001b[0m\n\u001b[1;32m      3\u001b[0m \u001b[0;32mimport\u001b[0m \u001b[0mos\u001b[0m\u001b[0;34m\u001b[0m\u001b[0;34m\u001b[0m\u001b[0m\n\u001b[1;32m      4\u001b[0m \u001b[0;32mimport\u001b[0m \u001b[0mmath\u001b[0m\u001b[0;34m\u001b[0m\u001b[0;34m\u001b[0m\u001b[0m\n\u001b[1;32m      5\u001b[0m \u001b[0;32mfrom\u001b[0m \u001b[0mpptx\u001b[0m \u001b[0;32mimport\u001b[0m \u001b[0mPresentation\u001b[0m\u001b[0;34m\u001b[0m\u001b[0;34m\u001b[0m\u001b[0m\n",
            "\u001b[0;31mModuleNotFoundError\u001b[0m: No module named 'pytube'",
            "",
            "\u001b[0;31m---------------------------------------------------------------------------\u001b[0;32m\nNOTE: If your import is failing due to a missing package, you can\nmanually install dependencies using either !pip or !apt.\n\nTo view examples of installing some common dependencies, click the\n\"Open Examples\" button below.\n\u001b[0;31m---------------------------------------------------------------------------\u001b[0m\n"
          ],
          "errorDetails": {
            "actions": [
              {
                "action": "open_url",
                "actionText": "Open Examples",
                "url": "/notebooks/snippets/importing_libraries.ipynb"
              }
            ]
          }
        }
      ]
    },
    {
      "cell_type": "markdown",
      "source": [
        "\n",
        "### Overall Software Documentation\n",
        "\n",
        "This software suite comprises three distinct functions, each addressing a specific aspect of media processing and presentation. The first function, `download_youtube_as_mp4`, allows users to download YouTube videos in MP4 format at a specified resolution, utilizing the `pytube` library. The second function, `extract_frames_cv2`, is designed for extracting frames from a given video at specified time intervals, leveraging the `cv2` (OpenCV) library for frame processing. Finally, the `images_to_ppt` function enables the creation of a PowerPoint presentation from a collection of images stored in a specified directory, using the `python-pptx` library. Together, these functions provide a comprehensive toolkit for video downloading, frame extraction, and presentation creation, useful in various multimedia and content creation workflows."
      ],
      "metadata": {
        "id": "hLb_v7YP8k1j"
      }
    },
    {
      "cell_type": "code",
      "source": [
        "### Documentation for Each Function\n",
        "\n",
        "#### 1. `download_youtube_as_mp4(video_url, resolution)`\n",
        "   #- **video_url** *(string)*: The URL of the YouTube video to be downloaded.\n",
        "   #- **resolution** *(string)*: The desired resolution of the downloaded video (e.g., '720p', '1080p').\n",
        "\n",
        "#### 2. `extract_frames_cv2(video_path, seconds)`\n",
        "   #- **video_path** *(string)*: The file path of the video from which frames are to be extracted.\n",
        "   #- **seconds** *(int/float)*: Interval in seconds at which frames will be extracted.\n",
        "\n",
        "#### 3. `images_to_ppt(directory, ppt_file)`\n",
        "   #- **directory** *(string)*: The path to the directory containing the images to be included in the PowerPoint presentation.\n",
        "   #- **ppt_file** *(string)*: The name or path for the output PowerPoint file.\n",
        "\n",
        "def download_youtube_as_mp4(video_url, resolution):\n",
        "\n",
        "\n",
        "    # Create a YouTube object and get the video stream\n",
        "    youtube = YouTube(video_url)\n",
        "    video_stream = youtube.streams.filter(progressive=True, file_extension='mp4', res=resolution).first()\n",
        "\n",
        "    if video_stream is None:\n",
        "        print(f\"No video found with resolution {resolution}\")\n",
        "        return None\n",
        "\n",
        "    # Function to track the progress of the download\n",
        "    def progress_function(stream, chunk, bytes_remaining):\n",
        "        current = (stream.filesize - bytes_remaining)\n",
        "        percent = (current / stream.filesize) * 100\n",
        "        print(f\"Downloaded {current} of {stream.filesize} bytes ({percent:.2f}%)\")\n",
        "\n",
        "    # Register the progress function\n",
        "    youtube.register_on_progress_callback(progress_function)\n",
        "\n",
        "    # Download the video stream as an MP4 file\n",
        "    mp4_file_path = video_stream.download()\n",
        "\n",
        "    return mp4_file_path\n"
      ],
      "metadata": {
        "id": "yMnDamOS5EF6"
      },
      "execution_count": 4,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "\n",
        "\n",
        "def images_to_ppt(directory, ppt_file):\n",
        "    # Append '.pptx' to the file name if it's not already there\n",
        "    if not ppt_file.endswith('.pptx'):\n",
        "        ppt_file += '.pptx'\n",
        "\n",
        "    prs = Presentation()  # Create a new presentation\n",
        "\n",
        "    # Get a list of all the image files in the directory\n",
        "    images = [f for f in os.listdir(directory) if f.endswith(('.png', '.jpg', '.jpeg', '.gif'))]\n",
        "\n",
        "     # Sort the list of image filenames\n",
        "    images = sorted(images)\n",
        "    # Add a slide with an image for each image file\n",
        "    for image_file in images:\n",
        "        slide_layout = prs.slide_layouts[5]  # Use the blank slide layout (layout 5)\n",
        "        slide = prs.slides.add_slide(slide_layout)\n",
        "\n",
        "        width = Inches(10)  # Define width and height of the image on the slide\n",
        "        height = Inches(7.5)\n",
        "        left = 0\n",
        "        top = 0\n",
        "\n",
        "        # Add the image\n",
        "        slide.shapes.add_picture(os.path.join(directory, image_file), left, top, width, height)\n",
        "\n",
        "    # Save the presentation\n",
        "    prs.save(ppt_file)\n",
        "\n",
        "\n"
      ],
      "metadata": {
        "id": "-NmbeushvjKh"
      },
      "execution_count": 17,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "\n",
        "def extract_frames_cv2(video_path, seconds):\n",
        "    vidcap = cv2.VideoCapture(video_path)\n",
        "    fps = vidcap.get(cv2.CAP_PROP_FPS)  # Gets the frames per second\n",
        "    multiplier = round(fps * seconds)\n",
        "\n",
        "    # Debug lines\n",
        "    print(f\"FPS: {fps}\")\n",
        "    print(f\"Multiplier: {multiplier}\")\n",
        "\n",
        "    directory_name = os.path.basename(video_path)[:30]\n",
        "    if not os.path.exists(directory_name):\n",
        "        os.makedirs(directory_name)\n",
        "\n",
        "    frame_number = 0\n",
        "    while vidcap.isOpened():\n",
        "\n",
        "        frame_id = int(round(vidcap.get(1)))  # Current frame number\n",
        "        ret, frame = vidcap.read()\n",
        "\n",
        "        # Debug line\n",
        "        #print(f\"Frame ID: {frame_id}, Return value: {ret}\")\n",
        "        #print(ret)\n",
        "        if not ret:\n",
        "            break\n",
        "        #print(frame)\n",
        "        if frame_id % multiplier == 0:\n",
        "            frame_number += 1\n",
        "            print('Has reaches the multiplier')\n",
        "            # Pad frame_number with leading zeros\n",
        "            padded_frame_number = str(frame_number).zfill(5)  # Pad with 5 zeros for example\n",
        "            cv2.imwrite(os.path.join(directory_name, f\"frame{padded_frame_number}.jpg\"), frame)\n",
        "\n",
        "    vidcap.release()\n",
        "    return directory_name\n",
        "\n",
        "# Example usage\n",
        "extract_frames_cv2(\"your_video_file_here.mp4\", 1)\n",
        "\n",
        "\n",
        "\n",
        "\n",
        "\n"
      ],
      "metadata": {
        "colab": {
          "base_uri": "https://localhost:8080/",
          "height": 71
        },
        "id": "e6wdSPwihI1k",
        "outputId": "8edfc779-94ab-41be-8f13-8c991dbc907c"
      },
      "execution_count": 6,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stdout",
          "text": [
            "FPS: 0.0\n",
            "Multiplier: 0\n"
          ]
        },
        {
          "output_type": "execute_result",
          "data": {
            "text/plain": [
              "'your_video_file_here.mp4'"
            ],
            "application/vnd.google.colaboratory.intrinsic+json": {
              "type": "string"
            }
          },
          "metadata": {},
          "execution_count": 6
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "import os\n",
        "import pickle\n",
        "from PIL import Image\n",
        "\n",
        "def pickle_images(directory):\n",
        "    images = []  # Create a list to store the images.\n",
        "\n",
        "    # Sort the list of image filenames\n",
        "    sorted_image_names = sorted(os.listdir(directory))\n",
        "\n",
        "    for image_name in sorted_image_names:\n",
        "        image_path = os.path.join(directory, image_name)  # Construct the full image path\n",
        "\n",
        "        try:\n",
        "            with Image.open(image_path) as image:  # Open the image\n",
        "                images.append(image.copy())  # Append a copy of the image to the list\n",
        "        except IOError:\n",
        "            print(f\"Unable to open {image_name}\")  # Handle the exception for an unreadable file\n",
        "\n",
        "    # Pickle the list of images.\n",
        "    with open('images.pkl', 'wb') as f:\n",
        "        pickle.dump(images, f)\n",
        "\n"
      ],
      "metadata": {
        "id": "NEzgeBl27a7i"
      },
      "execution_count": null,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "!jupyter nbconvert --to script youtube_to_images.ipynb"
      ],
      "metadata": {
        "colab": {
          "base_uri": "https://localhost:8080/"
        },
        "id": "DRTUu7Zj4XiL",
        "outputId": "89697b02-4475-4c31-b490-a9ae8ab63fdc"
      },
      "execution_count": 2,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stdout",
          "text": [
            "[NbConvertApp] WARNING | pattern 'youtube_to_images.ipynb' matched no files\n",
            "This application is used to convert notebook files (*.ipynb)\n",
            "        to various other formats.\n",
            "\n",
            "        WARNING: THE COMMANDLINE INTERFACE MAY CHANGE IN FUTURE RELEASES.\n",
            "\n",
            "Options\n",
            "=======\n",
            "The options below are convenience aliases to configurable class-options,\n",
            "as listed in the \"Equivalent to\" description-line of the aliases.\n",
            "To see all configurable class-options for some <cmd>, use:\n",
            "    <cmd> --help-all\n",
            "\n",
            "--debug\n",
            "    set log level to logging.DEBUG (maximize logging output)\n",
            "    Equivalent to: [--Application.log_level=10]\n",
            "--show-config\n",
            "    Show the application's configuration (human-readable format)\n",
            "    Equivalent to: [--Application.show_config=True]\n",
            "--show-config-json\n",
            "    Show the application's configuration (json format)\n",
            "    Equivalent to: [--Application.show_config_json=True]\n",
            "--generate-config\n",
            "    generate default config file\n",
            "    Equivalent to: [--JupyterApp.generate_config=True]\n",
            "-y\n",
            "    Answer yes to any questions instead of prompting.\n",
            "    Equivalent to: [--JupyterApp.answer_yes=True]\n",
            "--execute\n",
            "    Execute the notebook prior to export.\n",
            "    Equivalent to: [--ExecutePreprocessor.enabled=True]\n",
            "--allow-errors\n",
            "    Continue notebook execution even if one of the cells throws an error and include the error message in the cell output (the default behaviour is to abort conversion). This flag is only relevant if '--execute' was specified, too.\n",
            "    Equivalent to: [--ExecutePreprocessor.allow_errors=True]\n",
            "--stdin\n",
            "    read a single notebook file from stdin. Write the resulting notebook with default basename 'notebook.*'\n",
            "    Equivalent to: [--NbConvertApp.from_stdin=True]\n",
            "--stdout\n",
            "    Write notebook output to stdout instead of files.\n",
            "    Equivalent to: [--NbConvertApp.writer_class=StdoutWriter]\n",
            "--inplace\n",
            "    Run nbconvert in place, overwriting the existing notebook (only\n",
            "            relevant when converting to notebook format)\n",
            "    Equivalent to: [--NbConvertApp.use_output_suffix=False --NbConvertApp.export_format=notebook --FilesWriter.build_directory=]\n",
            "--clear-output\n",
            "    Clear output of current file and save in place,\n",
            "            overwriting the existing notebook.\n",
            "    Equivalent to: [--NbConvertApp.use_output_suffix=False --NbConvertApp.export_format=notebook --FilesWriter.build_directory= --ClearOutputPreprocessor.enabled=True]\n",
            "--no-prompt\n",
            "    Exclude input and output prompts from converted document.\n",
            "    Equivalent to: [--TemplateExporter.exclude_input_prompt=True --TemplateExporter.exclude_output_prompt=True]\n",
            "--no-input\n",
            "    Exclude input cells and output prompts from converted document.\n",
            "            This mode is ideal for generating code-free reports.\n",
            "    Equivalent to: [--TemplateExporter.exclude_output_prompt=True --TemplateExporter.exclude_input=True --TemplateExporter.exclude_input_prompt=True]\n",
            "--allow-chromium-download\n",
            "    Whether to allow downloading chromium if no suitable version is found on the system.\n",
            "    Equivalent to: [--WebPDFExporter.allow_chromium_download=True]\n",
            "--disable-chromium-sandbox\n",
            "    Disable chromium security sandbox when converting to PDF..\n",
            "    Equivalent to: [--WebPDFExporter.disable_sandbox=True]\n",
            "--show-input\n",
            "    Shows code input. This flag is only useful for dejavu users.\n",
            "    Equivalent to: [--TemplateExporter.exclude_input=False]\n",
            "--embed-images\n",
            "    Embed the images as base64 dataurls in the output. This flag is only useful for the HTML/WebPDF/Slides exports.\n",
            "    Equivalent to: [--HTMLExporter.embed_images=True]\n",
            "--sanitize-html\n",
            "    Whether the HTML in Markdown cells and cell outputs should be sanitized..\n",
            "    Equivalent to: [--HTMLExporter.sanitize_html=True]\n",
            "--log-level=<Enum>\n",
            "    Set the log level by value or name.\n",
            "    Choices: any of [0, 10, 20, 30, 40, 50, 'DEBUG', 'INFO', 'WARN', 'ERROR', 'CRITICAL']\n",
            "    Default: 30\n",
            "    Equivalent to: [--Application.log_level]\n",
            "--config=<Unicode>\n",
            "    Full path of a config file.\n",
            "    Default: ''\n",
            "    Equivalent to: [--JupyterApp.config_file]\n",
            "--to=<Unicode>\n",
            "    The export format to be used, either one of the built-in formats\n",
            "            ['asciidoc', 'custom', 'html', 'latex', 'markdown', 'notebook', 'pdf', 'python', 'rst', 'script', 'slides', 'webpdf']\n",
            "            or a dotted object name that represents the import path for an\n",
            "            ``Exporter`` class\n",
            "    Default: ''\n",
            "    Equivalent to: [--NbConvertApp.export_format]\n",
            "--template=<Unicode>\n",
            "    Name of the template to use\n",
            "    Default: ''\n",
            "    Equivalent to: [--TemplateExporter.template_name]\n",
            "--template-file=<Unicode>\n",
            "    Name of the template file to use\n",
            "    Default: None\n",
            "    Equivalent to: [--TemplateExporter.template_file]\n",
            "--theme=<Unicode>\n",
            "    Template specific theme(e.g. the name of a JupyterLab CSS theme distributed\n",
            "    as prebuilt extension for the lab template)\n",
            "    Default: 'light'\n",
            "    Equivalent to: [--HTMLExporter.theme]\n",
            "--sanitize_html=<Bool>\n",
            "    Whether the HTML in Markdown cells and cell outputs should be sanitized.This\n",
            "    should be set to True by nbviewer or similar tools.\n",
            "    Default: False\n",
            "    Equivalent to: [--HTMLExporter.sanitize_html]\n",
            "--writer=<DottedObjectName>\n",
            "    Writer class used to write the\n",
            "                                        results of the conversion\n",
            "    Default: 'FilesWriter'\n",
            "    Equivalent to: [--NbConvertApp.writer_class]\n",
            "--post=<DottedOrNone>\n",
            "    PostProcessor class used to write the\n",
            "                                        results of the conversion\n",
            "    Default: ''\n",
            "    Equivalent to: [--NbConvertApp.postprocessor_class]\n",
            "--output=<Unicode>\n",
            "    overwrite base name use for output files.\n",
            "                can only be used when converting one notebook at a time.\n",
            "    Default: ''\n",
            "    Equivalent to: [--NbConvertApp.output_base]\n",
            "--output-dir=<Unicode>\n",
            "    Directory to write output(s) to. Defaults\n",
            "                                  to output to the directory of each notebook. To recover\n",
            "                                  previous default behaviour (outputting to the current\n",
            "                                  working directory) use . as the flag value.\n",
            "    Default: ''\n",
            "    Equivalent to: [--FilesWriter.build_directory]\n",
            "--reveal-prefix=<Unicode>\n",
            "    The URL prefix for reveal.js (version 3.x).\n",
            "            This defaults to the reveal CDN, but can be any url pointing to a copy\n",
            "            of reveal.js.\n",
            "            For speaker notes to work, this must be a relative path to a local\n",
            "            copy of reveal.js: e.g., \"reveal.js\".\n",
            "            If a relative path is given, it must be a subdirectory of the\n",
            "            current directory (from which the server is run).\n",
            "            See the usage documentation\n",
            "            (https://nbconvert.readthedocs.io/en/latest/usage.html#reveal-js-html-slideshow)\n",
            "            for more details.\n",
            "    Default: ''\n",
            "    Equivalent to: [--SlidesExporter.reveal_url_prefix]\n",
            "--nbformat=<Enum>\n",
            "    The nbformat version to write.\n",
            "            Use this to downgrade notebooks.\n",
            "    Choices: any of [1, 2, 3, 4]\n",
            "    Default: 4\n",
            "    Equivalent to: [--NotebookExporter.nbformat_version]\n",
            "\n",
            "Examples\n",
            "--------\n",
            "\n",
            "    The simplest way to use nbconvert is\n",
            "\n",
            "            > jupyter nbconvert mynotebook.ipynb --to html\n",
            "\n",
            "            Options include ['asciidoc', 'custom', 'html', 'latex', 'markdown', 'notebook', 'pdf', 'python', 'rst', 'script', 'slides', 'webpdf'].\n",
            "\n",
            "            > jupyter nbconvert --to latex mynotebook.ipynb\n",
            "\n",
            "            Both HTML and LaTeX support multiple output templates. LaTeX includes\n",
            "            'base', 'article' and 'report'.  HTML includes 'basic', 'lab' and\n",
            "            'classic'. You can specify the flavor of the format used.\n",
            "\n",
            "            > jupyter nbconvert --to html --template lab mynotebook.ipynb\n",
            "\n",
            "            You can also pipe the output to stdout, rather than a file\n",
            "\n",
            "            > jupyter nbconvert mynotebook.ipynb --stdout\n",
            "\n",
            "            PDF is generated via latex\n",
            "\n",
            "            > jupyter nbconvert mynotebook.ipynb --to pdf\n",
            "\n",
            "            You can get (and serve) a Reveal.js-powered slideshow\n",
            "\n",
            "            > jupyter nbconvert myslides.ipynb --to slides --post serve\n",
            "\n",
            "            Multiple notebooks can be given at the command line in a couple of\n",
            "            different ways:\n",
            "\n",
            "            > jupyter nbconvert notebook*.ipynb\n",
            "            > jupyter nbconvert notebook1.ipynb notebook2.ipynb\n",
            "\n",
            "            or you can specify the notebooks list in a config file, containing::\n",
            "\n",
            "                c.NbConvertApp.notebooks = [\"my_notebook.ipynb\"]\n",
            "\n",
            "            > jupyter nbconvert --config mycfg.py\n",
            "\n",
            "To see all available configurables, use `--help-all`.\n",
            "\n"
          ]
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [],
      "metadata": {
        "id": "0kOef6FMhaA5"
      },
      "execution_count": 6,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [],
      "metadata": {
        "id": "zb_QtcDLmT0C"
      },
      "execution_count": 8,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "video_url = 'https://www.youtube.com/watch?v=PaCmpygFfXo&t=2657s&ab_channel=AndrejKarpathy'"
      ],
      "metadata": {
        "id": "oyUzADOTUCRJ"
      },
      "execution_count": 9,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "downloaded_video = download_youtube_as_mp4(video_url, '720p')"
      ],
      "metadata": {
        "id": "N2yig0Y15KKP",
        "colab": {
          "base_uri": "https://localhost:8080/"
        },
        "outputId": "9e951f35-8379-42b3-d1bb-06451d0ee862"
      },
      "execution_count": 10,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stdout",
          "text": [
            "Downloaded 9437184 of 196842631 bytes (4.79%)\n",
            "Downloaded 18874368 of 196842631 bytes (9.59%)\n",
            "Downloaded 28311552 of 196842631 bytes (14.38%)\n",
            "Downloaded 37748736 of 196842631 bytes (19.18%)\n",
            "Downloaded 47185920 of 196842631 bytes (23.97%)\n",
            "Downloaded 56623104 of 196842631 bytes (28.77%)\n",
            "Downloaded 66060288 of 196842631 bytes (33.56%)\n",
            "Downloaded 75497472 of 196842631 bytes (38.35%)\n",
            "Downloaded 84934656 of 196842631 bytes (43.15%)\n",
            "Downloaded 94371840 of 196842631 bytes (47.94%)\n",
            "Downloaded 103809024 of 196842631 bytes (52.74%)\n",
            "Downloaded 113246208 of 196842631 bytes (57.53%)\n",
            "Downloaded 122683392 of 196842631 bytes (62.33%)\n",
            "Downloaded 132120576 of 196842631 bytes (67.12%)\n",
            "Downloaded 141557760 of 196842631 bytes (71.91%)\n",
            "Downloaded 150994944 of 196842631 bytes (76.71%)\n",
            "Downloaded 160432128 of 196842631 bytes (81.50%)\n",
            "Downloaded 169869312 of 196842631 bytes (86.30%)\n",
            "Downloaded 179306496 of 196842631 bytes (91.09%)\n",
            "Downloaded 188743680 of 196842631 bytes (95.89%)\n",
            "Downloaded 196842631 of 196842631 bytes (100.00%)\n"
          ]
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "directory_name = extract_frames_cv2(downloaded_video, 30)"
      ],
      "metadata": {
        "id": "lWHXy8VVh98-",
        "colab": {
          "base_uri": "https://localhost:8080/"
        },
        "outputId": "bc54f16f-ea98-4a33-8b2e-1ae4005f5c81"
      },
      "execution_count": 11,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stdout",
          "text": [
            "FPS: 29.97002997002997\n",
            "Multiplier: 899\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n",
            "Has reaches the multiplier\n"
          ]
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "directory_name"
      ],
      "metadata": {
        "colab": {
          "base_uri": "https://localhost:8080/",
          "height": 35
        },
        "id": "jK-aEg5q_6j6",
        "outputId": "c8257699-4396-4790-9f9f-30e323654983"
      },
      "execution_count": 12,
      "outputs": [
        {
          "output_type": "execute_result",
          "data": {
            "text/plain": [
              "'The spelled-out intro to langu'"
            ],
            "application/vnd.google.colaboratory.intrinsic+json": {
              "type": "string"
            }
          },
          "metadata": {},
          "execution_count": 12
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "pickle_images(directory_name)"
      ],
      "metadata": {
        "id": "tJGyYq9e-yPs"
      },
      "execution_count": 13,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "images_to_ppt(directory_name, 'makemoreANDmore')"
      ],
      "metadata": {
        "id": "V2D6ZjaLtST-"
      },
      "execution_count": 18,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "%%capture\n"
      ],
      "metadata": {
        "id": "2Q5kuSkN2i6j",
        "colab": {
          "base_uri": "https://localhost:8080/"
        },
        "outputId": "ad4b3098-4595-4d6e-a0cf-dc8b32fedabe"
      },
      "execution_count": 15,
      "outputs": [
        {
          "output_type": "stream",
          "name": "stderr",
          "text": [
            "UsageError: %%capture is a cell magic, but the cell body is empty.\n"
          ]
        }
      ]
    },
    {
      "cell_type": "code",
      "source": [
        "\n"
      ],
      "metadata": {
        "id": "b9XP2FsJ39Cu"
      },
      "execution_count": null,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "'''from moviepy.editor import VideoFileClip\n",
        "\n",
        "def extract_frames_VideoFileClip(video_path, seconds):\n",
        "    clip = VideoFileClip(video_path)\n",
        "    fps = clip.fps  # Gets the frames per second\n",
        "\n",
        "    for i in range(0, int(clip.duration), seconds):\n",
        "        frame = clip.get_frame(i)\n",
        "        frame.save_frame(f\"frame{i}.jpg\", t=i)\n",
        "\n",
        "extract_frames('path_to_your_video.mp4', N)  # Replace N with the number of seconds\n",
        "'''\n"
      ],
      "metadata": {
        "id": "-H9eM9aXp8cx"
      },
      "execution_count": null,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "'''import imageio\n",
        "\n",
        "reader = imageio.get_reader_imageio('path_to_your_video.mp4')\n",
        "fps = reader.get_meta_data()['fps']  # Gets the frames per second\n",
        "\n",
        "for i, im in enumerate(reader):\n",
        "    if i % (fps * N) == 0:  # Replace N with the number of seconds\n",
        "        imageio.imwrite(f'frame{i}.jpg', im)'''"
      ],
      "metadata": {
        "id": "mYMrIdClrVtO"
      },
      "execution_count": null,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "print(var)"
      ],
      "metadata": {
        "id": "_UibiW7z2x2J"
      },
      "execution_count": null,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [],
      "metadata": {
        "id": "xji1uu01rn0s"
      },
      "execution_count": null,
      "outputs": []
    },
    {
      "cell_type": "code",
      "source": [
        "images_to_ppt('/content/Daniel Chen - Install Python Q' , 'Building an Audio Transcription App3.pptx')"
      ],
      "metadata": {
        "id": "0gWkSOludndL"
      },
      "execution_count": null,
      "outputs": []
    }
  ]
}